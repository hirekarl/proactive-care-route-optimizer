"""
build.py — Generate deck/presentation.pptx for the Proactive Care-Route Optimizer mini L2L.
Usage: uv run python build.py
"""
from __future__ import annotations

from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

# ---------------------------------------------------------------------------
# Paths
# ---------------------------------------------------------------------------

ROOT = Path(__file__).parent
LOGO_PATH = ROOT.parent / "frontend" / "public" / "care-route-logo.png"
OUT_PATH = ROOT / "presentation.pptx"

# ---------------------------------------------------------------------------
# Slide dimensions — widescreen 16:9
# ---------------------------------------------------------------------------

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# ---------------------------------------------------------------------------
# Color palette (matched from frontend design system)
# ---------------------------------------------------------------------------

BG = RGBColor(0x05, 0x05, 0x08)
WHITE = RGBColor(0xF8, 0xFA, 0xFC)
TEXT_SEC = RGBColor(0xCB, 0xD5, 0xE1)
TEXT_MUTED = RGBColor(0x94, 0xA3, 0xB8)
CARD_FILL = RGBColor(0x0F, 0x17, 0x2A)
CARD_BORDER = RGBColor(0x1E, 0x29, 0x3B)
SKY = RGBColor(0x38, 0xBD, 0xF8)
FUCHSIA = RGBColor(0xF0, 0xAB, 0xFC)
VIOLET = RGBColor(0xA7, 0x8B, 0xFA)
ROSE = RGBColor(0xFD, 0xA4, 0xAF)
EMERALD = RGBColor(0x34, 0xD3, 0x99)

FONT = "Segoe UI"

# Layout anchors
MARGIN_L = 0.65
LOGO_X, LOGO_Y = 0.65, 0.22
LOGO_SIZE = 0.44
CONTENT_TOP = 1.30

# ---------------------------------------------------------------------------
# Low-level helpers
# ---------------------------------------------------------------------------


def _rgb(color: RGBColor) -> str:
    return f"{color[0]:02X}{color[1]:02X}{color[2]:02X}"


def set_bg(slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def set_shape_alpha(shape, alpha_pct: int) -> None:
    """Inject a:alpha into solid fill XML (alpha_pct 0=invisible, 100=opaque)."""
    sp_pr = shape._element.spPr
    solid = sp_pr.find(qn("a:solidFill"))
    if solid is None:
        return
    srgb = solid.find(qn("a:srgbClr"))
    if srgb is None:
        return
    existing = srgb.find(qn("a:alpha"))
    if existing is not None:
        srgb.remove(existing)
    alpha_el = etree.SubElement(srgb, qn("a:alpha"))
    alpha_el.set("val", str(int(alpha_pct * 1000)))


def send_to_back(slide, shape) -> None:
    sp_tree = slide.shapes._spTree
    sp_tree.remove(shape._element)
    sp_tree.insert(2, shape._element)


def _set_para_spacing(tf, space_after_pt: float = 0, line_spacing: float | None = None) -> None:
    for para in tf.paragraphs:
        pPr = para._p.get_or_add_pPr()
        if space_after_pt:
            spcAft = etree.SubElement(pPr, qn("a:spcAft"))
            spcPct = etree.SubElement(spcAft, qn("a:spcPts"))
            spcPct.set("val", str(int(space_after_pt * 100)))
        if line_spacing is not None:
            lnSpc = etree.SubElement(pPr, qn("a:lnSpc"))
            spcPct = etree.SubElement(lnSpc, qn("a:spcPct"))
            spcPct.set("val", str(int(line_spacing * 100000)))


# ---------------------------------------------------------------------------
# Shape primitives
# ---------------------------------------------------------------------------


def add_glow(slide, x: float, y: float, w: float, h: float, color: RGBColor, alpha_pct: int):
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE as MSO

    shape = slide.shapes.add_shape(MSO.OVAL, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    set_shape_alpha(shape, alpha_pct)
    shape.line.fill.background()
    send_to_back(slide, shape)
    return shape


def add_rect(slide, x: float, y: float, w: float, h: float, fill: RGBColor, border: RGBColor | None = None):
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE as MSO

    shape = slide.shapes.add_shape(MSO.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if border:
        shape.line.color.rgb = border
        shape.line.width = Pt(0.75)
    else:
        shape.line.fill.background()
    return shape


def add_glass_card(slide, x: float, y: float, w: float, h: float, border_color: RGBColor | None = None):
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE as MSO

    shape = slide.shapes.add_shape(MSO.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD_FILL
    bc = border_color or CARD_BORDER
    shape.line.color.rgb = bc
    shape.line.width = Pt(0.75)
    shape.adjustments[0] = 0.05
    return shape


def add_severity_strip(slide, x: float, y: float, h: float, color: RGBColor):
    return add_rect(slide, x, y, 0.055, h, color)


def add_oval(slide, x: float, y: float, size: float, color: RGBColor):
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE as MSO

    shape = slide.shapes.add_shape(MSO.OVAL, Inches(x), Inches(y), Inches(size), Inches(size))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


# ---------------------------------------------------------------------------
# Text helpers
# ---------------------------------------------------------------------------


def _apply_run_font(run, font_name: str, size_pt: float, bold: bool, color: RGBColor) -> None:
    run.font.name = font_name
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = color


def add_textbox(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    font_size: float = 12,
    bold: bool = False,
    color: RGBColor = WHITE,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    word_wrap: bool = True,
    line_spacing: float | None = None,
):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = text
    _apply_run_font(run, FONT, font_size, bold, color)
    if line_spacing is not None:
        pPr = para._p.get_or_add_pPr()
        lnSpc = etree.SubElement(pPr, qn("a:lnSpc"))
        spcPct = etree.SubElement(lnSpc, qn("a:spcPct"))
        spcPct.set("val", str(int(line_spacing * 100000)))
    return tf


def add_textbox_lines(
    slide,
    lines: list[tuple[str, float, bool, RGBColor]],
    x: float,
    y: float,
    w: float,
    h: float,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    word_wrap: bool = True,
    line_spacing: float | None = None,
):
    """Multi-paragraph textbox. lines = [(text, font_size, bold, color), ...]"""
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    for i, (text, size, bold, color) in enumerate(lines):
        if i == 0:
            para = tf.paragraphs[0]
        else:
            para = tf.add_paragraph()
        para.alignment = align
        run = para.add_run()
        run.text = text
        _apply_run_font(run, FONT, size, bold, color)
        if line_spacing is not None:
            pPr = para._p.get_or_add_pPr()
            lnSpc = etree.SubElement(pPr, qn("a:lnSpc"))
            spcPct = etree.SubElement(lnSpc, qn("a:spcPct"))
            spcPct.set("val", str(int(line_spacing * 100000)))
    return tf


def add_eyebrow(slide, text: str, x: float, y: float, w: float, color: RGBColor):
    tf = add_textbox(slide, text.upper(), x, y, w, 0.25, font_size=8.5, bold=True, color=color)
    run = tf.paragraphs[0].runs[0]
    rPr = run._r.find(qn("a:rPr"))
    if rPr is not None:
        rPr.set("spc", "140")
    return tf


def add_pill_badge(slide, text: str, x: float, y: float, w: float, fill_color: RGBColor, text_color: RGBColor):
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE as MSO

    shape = slide.shapes.add_shape(MSO.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.30))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.adjustments[0] = 0.5
    tf = shape.text_frame
    tf.word_wrap = False
    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER
    run = para.add_run()
    run.text = text
    _apply_run_font(run, FONT, 8.5, True, text_color)
    rPr = run._r.find(qn("a:rPr"))
    if rPr is not None:
        rPr.set("spc", "100")
    # Vertical centering via margins
    tf.margin_top = Pt(4)
    tf.margin_bottom = Pt(4)
    tf.margin_left = Pt(10)
    tf.margin_right = Pt(10)
    return shape


# ---------------------------------------------------------------------------
# Chrome (logo + brand + rule) — on every slide
# ---------------------------------------------------------------------------


def add_chrome(slide) -> None:
    if LOGO_PATH.exists():
        slide.shapes.add_picture(str(LOGO_PATH), Inches(LOGO_X), Inches(LOGO_Y), Inches(LOGO_SIZE), Inches(LOGO_SIZE))

    add_textbox(
        slide,
        "Care-Route Optimizer",
        LOGO_X + LOGO_SIZE + 0.08,
        LOGO_Y + 0.09,
        2.5,
        0.3,
        font_size=9,
        color=TEXT_MUTED,
    )

    # Thin separator rule below logo
    rule_y = LOGO_Y + LOGO_SIZE + 0.14
    add_rect(slide, MARGIN_L, rule_y, 13.333 - MARGIN_L * 2, 0.007, CARD_BORDER)


def add_slide_number(slide, n: int) -> None:
    add_textbox(
        slide,
        f"{n} / 5",
        11.8,
        7.1,
        1.3,
        0.3,
        font_size=8.5,
        color=TEXT_MUTED,
        align=PP_ALIGN.RIGHT,
    )


# ---------------------------------------------------------------------------
# Slide 1 — Data Insight
# ---------------------------------------------------------------------------


def add_bar_row(slide, label: str, pct: float, color: RGBColor, x: float, y: float, bar_w: float) -> None:
    label_w = 0.35
    gap = 0.1
    track_x = x + label_w + gap
    track_w = bar_w - label_w - gap - 0.55
    bar_h = 0.14

    add_textbox(slide, label, x, y - 0.05, label_w, 0.25, font_size=9.5, color=TEXT_MUTED, align=PP_ALIGN.RIGHT)
    add_rect(slide, track_x, y, track_w, bar_h, CARD_BORDER)
    add_rect(slide, track_x, y, track_w * pct, bar_h, color)

    pct_label = f"{int(round(pct * 100))}%"
    add_textbox(slide, pct_label, track_x + track_w * pct + 0.06, y - 0.05, 0.45, 0.25, font_size=9.5, bold=True, color=color)


def build_slide_1(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)

    add_glow(slide, 8.0, -2.5, 7.0, 6.0, SKY, 11)
    add_glow(slide, -2.0, 4.5, 5.5, 4.5, FUCHSIA, 8)

    add_chrome(slide)

    # Left column
    add_eyebrow(slide, "Data Insight", MARGIN_L, CONTENT_TOP, 5.0, SKY)

    add_textbox(slide, "79%", 0.42, 1.52, 5.5, 1.55, font_size=90, bold=True, color=SKY)

    add_textbox(
        slide,
        "of DFTA providers sit within a quarter mile\nof a building with a chronic elevator-outage history",
        MARGIN_L, 3.18, 5.3, 1.1,
        font_size=15, color=WHITE, word_wrap=True, line_spacing=1.35,
    )

    add_rect(slide, MARGIN_L, 4.38, 1.1, 0.007, SKY)

    add_textbox(slide, "+20%", MARGIN_L, 4.48, 3.5, 0.9, font_size=44, bold=True, color=FUCHSIA)

    add_textbox(
        slide,
        "rise in elevator complaint rates\nduring heat weeks",
        MARGIN_L, 5.42, 4.8, 0.8,
        font_size=13, color=TEXT_SEC, word_wrap=True, line_spacing=1.3,
    )

    # Right glass card — borough bar chart
    card_x, card_y, card_w, card_h = 6.55, 1.25, 6.2, 5.8
    add_glass_card(slide, card_x, card_y, card_w, card_h)

    add_eyebrow(slide, "Building Risk Exposure by Borough", card_x + 0.3, card_y + 0.28, 5.6, TEXT_MUTED)

    bars = [
        ("BK", 0.82, ROSE),
        ("BX", 0.78, ROSE),
        ("MN", 0.71, FUCHSIA),
        ("QN", 0.65, SKY),
        ("SI", 0.41, EMERALD),
    ]
    bar_start_y = 2.10
    bar_gap = 0.72
    bar_area_x = card_x + 0.22
    bar_area_w = card_w - 0.28

    for i, (label, pct, color) in enumerate(bars):
        add_bar_row(slide, label, pct, color, bar_area_x, bar_start_y + i * bar_gap, bar_area_w)

    add_textbox(
        slide,
        "% of providers within ¼ mi of a chronic-offender building",
        card_x + 0.3, 6.62, 5.6, 0.3,
        font_size=7.5, color=TEXT_MUTED,
    )

    add_slide_number(slide, 1)


# ---------------------------------------------------------------------------
# Slide 2 — The Problem It Reveals
# ---------------------------------------------------------------------------


def build_slide_2(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)

    add_glow(slide, -2.5, 3.5, 6.5, 5.0, ROSE, 10)
    add_glow(slide, 8.5, -1.5, 5.5, 4.0, FUCHSIA, 7)

    add_chrome(slide)

    # Left column
    add_eyebrow(slide, "The Problem It Reveals", MARGIN_L, CONTENT_TOP, 5.5, ROSE)

    add_textbox_lines(
        slide,
        [
            ("When the elevator stops,", 26, True, WHITE),
            ("care stops too.", 26, True, ROSE),
        ],
        MARGIN_L, 1.55, 5.5, 1.1,
        line_spacing=1.2,
    )

    add_textbox(
        slide,
        "When those elevators break, homebound seniors on upper floors lose access to meals and care — and today's dispatchers don't find out until a worker is already stranded in the lobby.",
        MARGIN_L, 2.85, 5.4, 2.0,
        font_size=14.5, color=TEXT_SEC, word_wrap=True, line_spacing=1.45,
    )

    # Right alert card
    card_x, card_y, card_w, card_h = 6.75, 1.55, 6.0, 5.45
    add_glass_card(slide, card_x, card_y, card_w, card_h, ROSE)
    add_severity_strip(slide, card_x, card_y, card_h, ROSE)

    content_x = card_x + 0.28

    add_pill_badge(slide, "CRITICAL", content_x, card_y + 0.28, 1.0, ROSE, CARD_FILL)

    add_textbox(slide, "Maria C.  ·  Route 7B  ·  Stop 12", content_x, card_y + 0.72, 5.4, 0.3, font_size=12, bold=True, color=WHITE)
    add_textbox(slide, "247 Delancey St, Manhattan — Floor 9", content_x, card_y + 1.05, 5.4, 0.3, font_size=11, color=TEXT_SEC)

    add_rect(slide, content_x, card_y + 1.42, 5.3, 0.007, ROSE)

    add_textbox(slide, "ELEVATOR STATUS", content_x, card_y + 1.55, 2.5, 0.25, font_size=8, bold=True, color=TEXT_MUTED)
    add_textbox(slide, "Out of service  ·  Single unit  ·  3rd outage this month", content_x, card_y + 1.82, 5.3, 0.3, font_size=11, color=ROSE)

    add_rect(slide, content_x, card_y + 2.22, 5.3, 0.007, CARD_BORDER)

    add_textbox(slide, "DISPATCHER NOTIFIED", content_x, card_y + 2.35, 2.5, 0.25, font_size=8, bold=True, color=TEXT_MUTED)
    add_textbox(slide, "47 minutes after scheduled arrival", content_x, card_y + 2.62, 5.3, 0.3, font_size=12, bold=True, color=ROSE)

    # Footer inner card
    inner_card = add_glass_card(slide, content_x, card_y + 3.18, 5.3, 0.9, CARD_BORDER)
    add_textbox(
        slide,
        "→  Worker arrived at scene, found lobby locked. Senior on floor 9 missed meal delivery.",
        content_x + 0.18, card_y + 3.32, 4.9, 0.65,
        font_size=10.5, color=TEXT_SEC, word_wrap=True, line_spacing=1.3,
    )

    add_slide_number(slide, 2)


# ---------------------------------------------------------------------------
# Slide 3 — The Solution I Built
# ---------------------------------------------------------------------------


def add_step_card(slide, x: float, y: float, w: float, h: float, label: str, sublabel: str, dot_color: RGBColor) -> None:
    add_glass_card(slide, x, y, w, h)
    dot_size = 0.20
    dot_x = x + 0.25
    dot_y = y + (h - dot_size) / 2
    add_oval(slide, dot_x, dot_y, dot_size, dot_color)
    text_x = dot_x + dot_size + 0.18
    text_y_label = y + h / 2 - 0.22
    add_textbox(slide, label, text_x, text_y_label, w - dot_size - 0.65, 0.30, font_size=13, bold=True, color=WHITE)
    add_textbox(slide, sublabel, text_x, text_y_label + 0.30, w - dot_size - 0.65, 0.28, font_size=10, color=TEXT_MUTED)


def build_slide_3(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)

    add_glow(slide, 8.0, -2.0, 6.5, 5.5, VIOLET, 12)
    add_glow(slide, -2.0, 4.0, 5.5, 5.0, SKY, 8)

    add_chrome(slide)

    # Left column
    add_eyebrow(slide, "The Solution I Built", MARGIN_L, CONTENT_TOP, 5.5, VIOLET)

    add_textbox_lines(
        slide,
        [
            ("Route screening,", 26, True, WHITE),
            ("before dispatch.", 26, True, VIOLET),
        ],
        MARGIN_L, 1.55, 5.5, 1.1,
        line_spacing=1.2,
    )

    add_textbox(
        slide,
        "I built a route-screening tool that cross-references every planned delivery stop against the live NYC elevator-outage feed and a heat-adjusted risk model, flagging hazardous stops before a worker ever leaves the depot.",
        MARGIN_L, 2.85, 5.4, 2.2,
        font_size=14.5, color=TEXT_SEC, word_wrap=True, line_spacing=1.45,
    )

    # Right — 3-step flow diagram
    flow_x = 7.0
    flow_w = 5.75
    card_h = 1.45
    gap = 0.52

    steps = [
        ("Live NYC Elevator Outage Feed", "DOB open data, updated in real time", SKY),
        ("Heat-Adjusted Risk Model", "Complaint rate modifier × 7-day forecast", FUCHSIA),
        ("Hazardous Stop Flagged — Before Dispatch", "Worker never leaves the depot for a stranded stop", EMERALD),
    ]

    for i, (label, sublabel, color) in enumerate(steps):
        cy = CONTENT_TOP + i * (card_h + gap)
        add_step_card(slide, flow_x, cy, flow_w, card_h, label, sublabel, color)
        # Connector between cards
        if i < 2:
            connector_color = steps[i][2]
            connector_y = cy + card_h
            connector_x = flow_x + flow_w / 2 - 0.01
            add_rect(slide, connector_x, connector_y, 0.02, gap, connector_color)

    add_slide_number(slide, 3)


# ---------------------------------------------------------------------------
# Slide 4 — What It Does
# ---------------------------------------------------------------------------


def add_feature_card(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    label: str,
    sublabel: str,
    color: RGBColor,
    dark_tint: RGBColor,
) -> None:
    add_glass_card(slide, x, y, w, h)
    dot_size = 0.18
    add_oval(slide, x + 0.22, y + 0.22, dot_size, color)
    add_eyebrow(slide, sublabel, x + 0.22, y + 0.50, w - 0.3, color)
    add_textbox(slide, label, x + 0.22, y + 0.78, w - 0.3, 0.8, font_size=13, bold=True, color=WHITE, word_wrap=True)
    # Dark tint strip at bottom
    add_rect(slide, x, y + h - 0.28, w, 0.28, dark_tint)


def build_slide_4(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)

    add_glow(slide, 8.5, -1.5, 6.0, 5.0, SKY, 11)
    add_glow(slide, -1.5, 4.5, 5.5, 4.0, VIOLET, 10)

    add_chrome(slide)

    # Left column
    add_eyebrow(slide, "What It Does", MARGIN_L, CONTENT_TOP, 5.5, SKY)

    add_textbox_lines(
        slide,
        [
            ("The Dispatcher", 26, True, WHITE),
            ("Dashboard", 26, True, SKY),
        ],
        MARGIN_L, 1.55, 5.5, 1.1,
        line_spacing=1.2,
    )

    add_textbox(
        slide,
        "The dispatcher dashboard surfaces real-time outage alerts, composite building risk scores, and a 7-day heat forecast — giving dispatchers the information to reroute proactively instead of scrambling reactively.",
        MARGIN_L, 2.85, 5.4, 2.2,
        font_size=14.5, color=TEXT_SEC, word_wrap=True, line_spacing=1.45,
    )

    # Right 2×2 feature card grid
    grid_x = 7.0
    grid_y = 1.35
    card_w = 2.85
    card_h = 2.75
    gap = 0.18

    features = [
        ("Real-Time\nOutage Alerts", "LIVE", ROSE, RGBColor(0x1A, 0x0B, 0x0E)),
        ("Building\nRisk Score", "COMPOSITE", SKY, RGBColor(0x07, 0x15, 0x20)),
        ("7-Day Heat\nForecast", "FORECAST", FUCHSIA, RGBColor(0x17, 0x09, 0x1A)),
        ("Proactive\nRerouting", "OUTCOME", EMERALD, RGBColor(0x07, 0x15, 0x10)),
    ]

    positions = [
        (grid_x, grid_y),
        (grid_x + card_w + gap, grid_y),
        (grid_x, grid_y + card_h + gap),
        (grid_x + card_w + gap, grid_y + card_h + gap),
    ]

    for (label, sublabel, color, tint), (fx, fy) in zip(features, positions):
        add_feature_card(slide, fx, fy, card_w, card_h, label, sublabel, color, tint)

    add_slide_number(slide, 4)


# ---------------------------------------------------------------------------
# Slide 5 — Why It Matters
# ---------------------------------------------------------------------------


def build_slide_5(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)

    add_glow(slide, -2.0, 3.0, 6.5, 6.0, EMERALD, 9)
    add_glow(slide, 9.0, -2.0, 6.0, 5.5, SKY, 11)

    add_chrome(slide)

    # Left column
    add_eyebrow(slide, "Why It Matters", MARGIN_L, CONTENT_TOP, 5.0, EMERALD)

    add_textbox(slide, "1.8M", 0.42, 1.52, 7.5, 1.6, font_size=96, bold=True, color=SKY)

    add_textbox(
        slide,
        "older New Yorkers depend on DFTA-contracted providers for meals and care.",
        MARGIN_L, 3.52, 7.5, 0.9,
        font_size=19, color=WHITE, word_wrap=True, line_spacing=1.3,
    )

    add_textbox(
        slide,
        "This is the first tool that shifts senior-care delivery from reactive to proactive.",
        MARGIN_L, 4.58, 8.2, 0.85,
        font_size=15, color=TEXT_SEC, word_wrap=True, line_spacing=1.35,
    )

    # Right — logo showcase
    if LOGO_PATH.exists():
        slide.shapes.add_picture(str(LOGO_PATH), Inches(10.2), Inches(1.38), Inches(2.6), Inches(2.6))

    add_textbox_lines(
        slide,
        [
            ("Care-Route", 17, True, WHITE),
            ("Optimizer", 13, False, TEXT_MUTED),
        ],
        9.9, 4.18, 3.2, 0.75,
        align=PP_ALIGN.CENTER,
        line_spacing=1.2,
    )

    # CTA pill — SKY → VIOLET gradient
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE as MSO

    pill = slide.shapes.add_shape(MSO.ROUNDED_RECTANGLE, Inches(2.2), Inches(6.28), Inches(4.2), Inches(0.68))
    pill.adjustments[0] = 0.5
    fill = pill.fill
    fill.gradient()
    gs = fill.gradient_stops
    gs[0].position = 0.0
    gs[0].color.rgb = SKY
    gs[1].position = 1.0
    gs[1].color.rgb = VIOLET
    fill.gradient_angle = 135.0
    pill.line.fill.background()

    tf = pill.text_frame
    tf.word_wrap = False
    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER
    run = para.add_run()
    run.text = "Reactive  →  Proactive"
    _apply_run_font(run, FONT, 14, True, CARD_FILL)
    tf.margin_top = Pt(6)

    add_slide_number(slide, 5)


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------


def build() -> None:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    for n, builder in enumerate(
        [build_slide_1, build_slide_2, build_slide_3, build_slide_4, build_slide_5],
        start=1,
    ):
        builder(prs)

    prs.save(OUT_PATH)
    print(f"Saved → {OUT_PATH}")


if __name__ == "__main__":
    build()
